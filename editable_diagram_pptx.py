#!/usr/bin/env python3
"""
Generate a fully-editable PowerPoint diagram (shapes + connectors).

This script recreates the high-level structure of the provided "encoder/transmission/decoder"
block diagram as PPTX shapes (NOT an embedded image).

Optionally, you can provide a reference image; it will be added on a separate slide for visual
comparison/alignment.

Usage:
  python3 editable_diagram_pptx.py -o /home/mazya/diagram_editable.pptx
  python3 editable_diagram_pptx.py -o out.pptx --reference-image /path/to/original.png
"""

from __future__ import annotations

import argparse
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Create an editable PPTX diagram using shapes.")
    p.add_argument("-o", "--output", required=True, help="Output .pptx path.")
    p.add_argument("--reference-image", default=None, help="Optional reference image path (png/jpg).")
    p.add_argument("--layout", choices=("widescreen", "standard"), default="widescreen")
    return p.parse_args()


def _rgb(hex6: str) -> RGBColor:
    hex6 = hex6.strip().lstrip("#")
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def _set_line(shape, color: RGBColor, width_pt: float = 1.5):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _set_fill(shape, color: RGBColor, transparency: float | None = None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        # python-pptx supports fill transparency on some versions; keep best-effort.
        try:
            shape.fill.fore_color.transparency = float(transparency)
        except Exception:
            pass


def _add_label_box(slide, x, y, w, h, text, fill, line, font_size=14, bold=False):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    _set_fill(box, fill)
    _set_line(box, line, 1.0)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = bold
    return box


def _add_title(slide, x, y, w, h, text):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    return tb


def _add_module(slide, x, y, w, h, title, inner_label=None):
    # Outer rounded module
    module = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    _set_fill(module, _rgb("DDEED8"))
    _set_line(module, _rgb("6AA84F"), 2.0)

    _add_title(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.35), title)

    # Inner "NN" placeholder (three stacked columns)
    pad = Inches(0.18)
    inner_x = x + pad
    inner_y = y + Inches(0.55)
    inner_w = w - 2 * pad
    inner_h = h - Inches(0.85)

    col_gap = Inches(0.15)
    col_w = (inner_w - 2 * col_gap) // 3
    for i in range(3):
        col = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            inner_x + i * (col_w + col_gap),
            inner_y,
            col_w,
            inner_h,
        )
        _set_fill(col, _rgb("FFFFFF"))
        _set_line(col, _rgb("6AA84F"), 1.5)
        # Add small "neurons"
        n = 5
        dot_r = Inches(0.08)
        for j in range(n):
            dy = inner_h * (j + 1) / (n + 1)
            dot = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                col.left + int(col_w / 2 - dot_r / 2),
                col.top + int(dy - dot_r / 2),
                dot_r,
                dot_r,
            )
            _set_fill(dot, _rgb("FFFFFF"))
            _set_line(dot, _rgb("000000"), 1.0)

    if inner_label:
        _add_label_box(
            slide,
            x + w - Inches(0.85),
            y + Inches(0.45),
            Inches(0.6),
            Inches(0.35),
            inner_label,
            _rgb("6AA84F"),
            _rgb("6AA84F"),
            font_size=14,
            bold=True,
        )

    return module


def _add_line_with_arrow(slide, x1, y1, x2, y2, color, width_pt=2.0, arrow_size_in=0.18):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    _set_line(line, color, width_pt)

    # Arrowhead as triangle shape (because this python-pptx build has no arrowhead enums)
    tri = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        x2 - Inches(arrow_size_in) // 2,
        y2 - Inches(arrow_size_in) // 2,
        Inches(arrow_size_in),
        Inches(arrow_size_in),
    )
    _set_fill(tri, color)
    tri.line.fill.background()  # no outline

    ang = math.degrees(math.atan2(float(y2 - y1), float(x2 - x1)))
    tri.rotation = ang + 90  # default triangle points up
    return line, tri


def build(prs: Presentation, reference_image: Path | None):
    # Slide size presets
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Reference (always present)
    slide_ref = prs.slides.add_slide(prs.slide_layouts[6])
    if reference_image:
        slide_ref.shapes.add_picture(str(reference_image), Inches(0), Inches(0), width=prs.slide_width)
    else:
        tb = slide_ref.shapes.add_textbox(Inches(0.75), Inches(1.75), prs.slide_width - Inches(1.5), Inches(1.5))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = "Reference slide\n\nProvide --reference-image /path/to/original.png to embed the original figure here."
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.bold = True

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Palette (approximate to original)
    green = _rgb("6AA84F")
    light_green = _rgb("DDEED8")
    gray = _rgb("EFEFEF")
    dark_gray = _rgb("666666")
    blue = _rgb("3C78D8")
    orange = _rgb("F6B26B")
    purple = _rgb("B4A7D6")

    # Modules
    enc = _add_module(slide, Inches(1.0), Inches(1.1), Inches(3.2), Inches(2.6), "Encoder (TX)", inner_label="wₑ")
    dec = _add_module(slide, Inches(9.2), Inches(1.1), Inches(3.2), Inches(2.6), "Decoder (RX)", inner_label="w_d")
    ctrl = _add_module(slide, Inches(1.0), Inches(4.35), Inches(4.2), Inches(2.6), "Metasurface Controller", inner_label="w_m")
    conf = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(4.35), Inches(4.4), Inches(2.6))
    _set_fill(conf, light_green)
    _set_line(conf, green, 2.0)
    _add_title(slide, Inches(9.1), Inches(4.45), Inches(4.2), Inches(0.55), "Trainable Metasurface\nConfiguration")

    # Transmission block
    tx = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(2.25), Inches(3.2), Inches(1.35))
    _set_fill(tx, gray)
    _set_line(tx, dark_gray, 2.0)
    tf = tx.text_frame
    tf.text = "Transmission\nT(ℋ(t), ϕ(t), s(t))"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)

    # Loss function block
    loss = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.35), Inches(0.35), Inches(2.6), Inches(0.75))
    _set_fill(loss, _rgb("FFFFFF"))
    _set_line(loss, _rgb("999999"), 2.0)
    tf = loss.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = "Loss Function"
    p0.alignment = PP_ALIGN.CENTER
    p0.font.size = Pt(16)
    p0.font.bold = True
    p1 = tf.add_paragraph()
    p1.text = "ℒ(o(t), ô(t))"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(14)

    # Signal boxes
    o_box = _add_label_box(slide, Inches(0.2), Inches(1.45), Inches(0.7), Inches(0.35), "o(t)", orange, _rgb("C27C0E"), 14, True)
    x_box = _add_label_box(slide, Inches(0.35), Inches(2.05), Inches(0.7), Inches(0.35), "x(t)", _rgb("FFF2CC"), _rgb("B7B7B7"), 14)
    h_box = _add_label_box(slide, Inches(0.35), Inches(2.55), Inches(0.7), Inches(0.35), "ℋ(t)", _rgb("FFF2CC"), _rgb("B7B7B7"), 14)

    s_box = _add_label_box(slide, Inches(4.35), Inches(2.6), Inches(0.75), Inches(0.35), "s(t)", purple, _rgb("674EA7"), 14, True)
    y_box = _add_label_box(slide, Inches(8.25), Inches(2.45), Inches(0.75), Inches(0.35), "y(t)", _rgb("EA9999"), _rgb("CC0000"), 14, True)
    n_box = _add_label_box(slide, Inches(8.25), Inches(2.9), Inches(0.75), Inches(0.35), "n(t)", _rgb("CFE2F3"), _rgb("3C78D8"), 14, True)
    ohat_box = _add_label_box(slide, Inches(12.5), Inches(2.35), Inches(0.75), Inches(0.35), "ô(t)", purple, _rgb("674EA7"), 14, True)

    phi_box = _add_label_box(slide, Inches(6.25), Inches(4.05), Inches(0.85), Inches(0.4), "ϕ(t)", purple, _rgb("674EA7"), 16, True)

    # Control type diamond and configuration mini-chain
    diamond = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(5.45), Inches(5.1), Inches(2.3), Inches(1.1))
    _set_fill(diamond, _rgb("FFFFFF"))
    _set_line(diamond, _rgb("000000"), 1.5)
    tf = diamond.text_frame
    tf.text = "Type of\nControl"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(14)

    # Trainable metasurface config chain (ω -> exp -> φbar)
    w_box = _add_label_box(slide, Inches(10.6), Inches(5.15), Inches(0.6), Inches(0.35), "ω", _rgb("6AA84F"), _rgb("38761D"), 16, True)
    exp_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.05), Inches(5.65), Inches(1.7), Inches(0.45))
    _set_fill(exp_box, _rgb("EFEFEF"))
    _set_line(exp_box, _rgb("999999"), 1.5)
    exp_box.text_frame.text = "exp(−j2π ω)"
    exp_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    exp_box.text_frame.paragraphs[0].font.size = Pt(13)
    phibar_box = _add_label_box(slide, Inches(10.6), Inches(6.25), Inches(0.7), Inches(0.35), "ϕ̄", purple, _rgb("674EA7"), 16, True)

    # Primary black flows
    _add_line_with_arrow(slide, o_box.left + o_box.width, o_box.top + o_box.height // 2, enc.left, enc.top + Inches(0.75), _rgb("000000"), 2.0)
    _add_line_with_arrow(slide, x_box.left + x_box.width, x_box.top + x_box.height // 2, enc.left, enc.top + Inches(1.25), _rgb("000000"), 2.0)
    _add_line_with_arrow(slide, h_box.left + h_box.width, h_box.top + h_box.height // 2, enc.left, enc.top + Inches(1.75), _rgb("000000"), 2.0)

    _add_line_with_arrow(slide, enc.left + enc.width, s_box.top + s_box.height // 2, s_box.left, s_box.top + s_box.height // 2, _rgb("000000"), 2.0)
    _add_line_with_arrow(slide, s_box.left + s_box.width, s_box.top + s_box.height // 2, tx.left, tx.top + tx.height // 2, _rgb("000000"), 2.0)

    _add_line_with_arrow(slide, tx.left + tx.width, tx.top + tx.height // 2, y_box.left, y_box.top + y_box.height // 2, _rgb("000000"), 2.0)
    _add_line_with_arrow(slide, y_box.left + y_box.width, y_box.top + y_box.height // 2, dec.left, dec.top + Inches(1.0), _rgb("000000"), 2.0)
    _add_line_with_arrow(slide, n_box.left + n_box.width, n_box.top + n_box.height // 2, dec.left, dec.top + Inches(1.5), _rgb("000000"), 2.0)
    _add_line_with_arrow(slide, dec.left + dec.width, ohat_box.left, dec.top + Inches(1.2), ohat_box.top + ohat_box.height // 2, _rgb("000000"), 2.0)

    # Control to transmission
    _add_line_with_arrow(slide, ctrl.left + ctrl.width, ctrl.top + Inches(1.4), phi_box.left, phi_box.top + phi_box.height // 2, _rgb("000000"), 2.0)
    _add_line_with_arrow(slide, phi_box.left + phi_box.width, phi_box.top + phi_box.height // 2, tx.left + tx.width // 2, tx.top + tx.height, _rgb("000000"), 2.0)
    _add_line_with_arrow(slide, h_box.left + h_box.width, h_box.top + h_box.height // 2, ctrl.left, ctrl.top + Inches(1.2), _rgb("000000"), 2.0)

    # Blue "learning" arrows (simplified)
    _add_line_with_arrow(slide, loss.left + loss.width // 2, loss.top + loss.height, enc.left + enc.width // 2, enc.top, blue, 2.0)
    _add_line_with_arrow(slide, loss.left + loss.width // 2, loss.top + loss.height, dec.left + dec.width // 2, dec.top, blue, 2.0)
    _add_line_with_arrow(slide, loss.left + loss.width // 2, loss.top + loss.height, ctrl.left + ctrl.width // 2, ctrl.top, blue, 2.0)
    _add_line_with_arrow(slide, loss.left + loss.width // 2, loss.top + loss.height, conf.left + conf.width // 2, conf.top, blue, 2.0)

    # "Fixed / reconfigurable" callouts (dashed lines)
    left_lbl = slide.shapes.add_textbox(Inches(4.0), Inches(6.45), Inches(2.2), Inches(0.4))
    left_lbl.text_frame.text = "Reconfigurable"
    left_lbl.text_frame.paragraphs[0].font.size = Pt(12)
    left_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    right_lbl = slide.shapes.add_textbox(Inches(7.1), Inches(6.45), Inches(2.2), Inches(0.4))
    right_lbl.text_frame.text = "Fixed\nConfiguration"
    right_lbl.text_frame.paragraphs[0].font.size = Pt(12)
    right_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Arrow from diamond to ctrl/conf (dashed)
    for x2, y2 in [(ctrl.left + ctrl.width // 2, ctrl.top), (conf.left + conf.width // 2, conf.top)]:
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            diamond.left + diamond.width // 2,
            diamond.top + diamond.height,
            x2,
            y2,
        )
        _set_line(line, _rgb("000000"), 1.25)
        try:
            line.line.dash_style = 2  # best-effort dash; enum may vary
        except Exception:
            pass


def main() -> int:
    args = _args()
    out = Path(args.output).expanduser().resolve()
    ref = Path(args.reference_image).expanduser().resolve() if args.reference_image else None
    if ref and not ref.exists():
        raise SystemExit(f"Reference image not found: {ref}")

    prs = Presentation()
    build(prs, ref)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote: {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
