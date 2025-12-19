#!/usr/bin/env python3
"""
Create a PowerPoint (.pptx) from a single image (e.g., a diagram).

Usage:
  python3 image_to_pptx.py /path/to/diagram.png -o diagram.pptx --title "My Diagram"
"""

from __future__ import annotations

import argparse
import os
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Convert a single image to a PowerPoint slide.")
    p.add_argument("image", help="Path to the input image (png/jpg/webp/etc).")
    p.add_argument(
        "-o",
        "--output",
        default=None,
        help="Output .pptx path (default: same name as input, .pptx).",
    )
    p.add_argument(
        "--layout",
        choices=("widescreen", "standard"),
        default="widescreen",
        help="Slide aspect ratio preset.",
    )
    p.add_argument(
        "--margin",
        type=float,
        default=0.35,
        help="Margin in inches around the image (default: 0.35).",
    )
    p.add_argument(
        "--title",
        default=None,
        help="Optional title text to add at the top of the slide.",
    )
    return p.parse_args()


def main() -> int:
    args = _parse_args()
    img_path = Path(args.image).expanduser().resolve()
    if not img_path.exists():
        raise SystemExit(f"Image not found: {img_path}")

    out_path = Path(args.output).expanduser().resolve() if args.output else img_path.with_suffix(".pptx")
    if out_path.suffix.lower() != ".pptx":
        out_path = out_path.with_suffix(".pptx")

    prs = Presentation()

    # Slide size presets
    if args.layout == "widescreen":
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)  # 4:3
        prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    margin = Inches(args.margin)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    title_h = Inches(0)
    if args.title:
        title_h = Inches(0.6)
        title_box = slide.shapes.add_textbox(margin, Inches(0.1), slide_w - 2 * margin, title_h)
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = args.title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.bold = True

    with Image.open(img_path) as im:
        img_w_px, img_h_px = im.size

    # Available space for the image
    top_reserved = margin + title_h
    avail_w = slide_w - 2 * margin
    avail_h = slide_h - top_reserved - margin

    # Scale image to fit while preserving aspect ratio
    scale = min(avail_w / img_w_px, avail_h / img_h_px)
    pic_w = int(img_w_px * scale)
    pic_h = int(img_h_px * scale)

    left = int((slide_w - pic_w) / 2)
    top = int(top_reserved + (avail_h - pic_h) / 2)

    slide.shapes.add_picture(str(img_path), left, top, width=pic_w, height=pic_h)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
